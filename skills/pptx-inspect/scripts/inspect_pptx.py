"""pptx(PowerPoint)ファイルのスライド構造（shape単位）を抽出してJSONで出力する。

pptx-tools スキルの実行スクリプト（progressive disclosure 第3段階）。
read_pptx.py が人間向けの要約（title/texts/tables/notes）を返すのに対し、
このスクリプトは edit_pptx.py で編集対象を指定するために必要な
shape単位の構造情報（shape_index・placeholder情報・表/画像の有無）を返す。

run_script ツールから
    python inspect_pptx.py <pptx_path> [--start-slide N] [--max-slides N]
の形で呼ばれる。read_pdf.py / read_pptx.py と同じページング設計を流用する。
"""

from __future__ import annotations

import argparse
import json
import sys
import traceback
from pathlib import Path

from _common import check_shape_overflow, describe_shape, setup_utf8_stdio, summarize_result, write_json_result

from pptx import Presentation
from pptx.exc import PackageNotFoundError


def extract_slide(slide) -> dict:
    layout = slide.slide_layout
    # layout_index は layout が属する slide_master 内での連番（prs.slide_layouts は
    # 最初の master のものしか見えないため、layout.slide_master 経由で正しい master から引く）。
    layout_index = list(layout.slide_master.slide_layouts).index(layout)
    shapes = [describe_shape(shape, i) for i, shape in enumerate(slide.shapes)]
    notes_present = False
    if slide.has_notes_slide:
        notes_text = slide.notes_slide.notes_text_frame.text
        notes_present = bool(notes_text.strip())
    return {
        "layout_name": layout.name,
        "layout_index": layout_index,
        "shapes": shapes,
        "notes_present": notes_present,
    }


def main() -> int:
    setup_utf8_stdio()
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx_path")
    parser.add_argument("--start-slide", type=int, default=1)
    parser.add_argument("--max-slides", type=int, default=20)
    args = parser.parse_args()

    path = Path(args.pptx_path)
    if not path.exists():
        print(f"ファイルが見つかりません: {args.pptx_path}", file=sys.stderr)
        return 1
    if path.is_dir():
        print(f"指定パスはディレクトリです（ファイル専用）: {args.pptx_path}", file=sys.stderr)
        return 1

    try:
        prs = Presentation(str(path))
    except PackageNotFoundError:
        print(f"pptxとして読み込めませんでした（壊れているか非対応形式の可能性）: {args.pptx_path}", file=sys.stderr)
        return 1
    except Exception as e:  # noqa: BLE001 - python-pptx側の想定外エラーもエラー終了に統一する
        print(f"pptxの読み込みに失敗しました: {e}\n{traceback.format_exc()}", file=sys.stderr)
        return 1

    total_slides = len(prs.slides)

    # スライドサイズをcm単位で取得
    from _common import _length_cm
    slide_width_cm = _length_cm(prs.slide_width)
    slide_height_cm = _length_cm(prs.slide_height)

    start_idx = max(args.start_slide, 1) - 1
    end_idx = min(start_idx + max(args.max_slides, 1), total_slides)

    slides_out = []
    warnings = []
    for i in range(start_idx, end_idx):
        slide = prs.slides[i]
        info = extract_slide(slide)
        slide_entry = {"index": i + 1, **info}
        slides_out.append(slide_entry)

        # 各shapeの境界はみ出しをチェック
        for shape_info in info.get("shapes", []):
            overflow_msg = check_shape_overflow(shape_info, slide_width_cm, slide_height_cm)
            if overflow_msg:
                warnings.append(f"スライド{i + 1}の{overflow_msg}")

    result = {
        "path": str(path),
        "total_slides": total_slides,
        "start_slide": start_idx + 1 if slides_out else None,
        "end_slide": end_idx if slides_out else None,
        "slide_width_cm": slide_width_cm,
        "slide_height_cm": slide_height_cm,
        "slides": slides_out,
    }
    if warnings:
        result["warnings"] = warnings
    summary = summarize_result(result, ["slides"])
    summary.update(write_json_result(result, "pptx_inspect", path))
    print(json.dumps(summary, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    sys.exit(main())
