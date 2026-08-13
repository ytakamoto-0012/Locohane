"""pptx-tools スキル内の各スクリプトが共有するヘルパー関数。

run_script からは直接実行されない（scripts/ prefix チェックを通らないため）。
scripts/ 配下の各スクリプトが同一ディレクトリから import して使う。
"""

from __future__ import annotations

import hashlib
import json
import os
import shutil
import sys
from datetime import datetime
from pathlib import Path

from pptx.enum.shapes import MSO_SHAPE_TYPE


def setup_utf8_stdio() -> None:
    """標準出力/標準エラーをUTF-8に固定する。

    Windows環境ではパイプ経由のPython子プロセスの既定エンコーディングが
    システムのANSIコードページ（例: cp932）になり、run_script側の
    encoding="utf-8" デコードと食い違って日本語が文字化けするため、
    各スクリプトの先頭で必ず呼ぶこと。
    """
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")


def backup_before_overwrite(path: Path) -> Path | None:
    """path が既に存在する場合、上書き直前に同じフォルダへタイムスタンプ付きで
    コピーしてバックアップを作成する。存在しなければ何もせず None を返す
    （新規作成時はバックアップ不要）。
    """
    if not path.exists():
        return None
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    backup_path = path.with_name(f"{path.stem}.bak_{timestamp}{path.suffix}")
    suffix_n = 2
    while backup_path.exists():
        backup_path = path.with_name(f"{path.stem}.bak_{timestamp}_{suffix_n}{path.suffix}")
        suffix_n += 1
    shutil.copy2(path, backup_path)
    return backup_path


def _length_cm(value) -> float | None:
    """python-pptxのLength（EMU）をcm単位のfloatへ変換する（None時はNoneのまま）。"""
    return round(value.cm, 2) if value is not None else None


def check_shape_overflow(shape_info: dict, slide_width_cm: float | None, slide_height_cm: float | None) -> str | None:
    """shapeがスライド境界をはみ出しているかをチェックする。

    はみ出していれば説明文を返す。はみ出していなければNoneを返す。
    微小許容誤差（0.05cm）を加味する。
    """
    if slide_width_cm is None or slide_height_cm is None:
        return None

    left_cm = shape_info.get("left_cm")
    top_cm = shape_info.get("top_cm")
    width_cm = shape_info.get("width_cm")
    height_cm = shape_info.get("height_cm")

    if left_cm is None or top_cm is None or width_cm is None or height_cm is None:
        return None

    tolerance = 0.05
    issues = []

    if left_cm < -tolerance:
        issues.append(f"左端を{abs(left_cm):.2f}cmはみ出し")
    if top_cm < -tolerance:
        issues.append(f"上端を{abs(top_cm):.2f}cmはみ出し")
    if left_cm + width_cm > slide_width_cm + tolerance:
        overshoot = left_cm + width_cm - slide_width_cm
        issues.append(f"右端を{overshoot:.2f}cmはみ出し")
    if top_cm + height_cm > slide_height_cm + tolerance:
        overshoot = top_cm + height_cm - slide_height_cm
        issues.append(f"下端を{overshoot:.2f}cmはみ出し")

    if not issues:
        return None

    shape_name = shape_info.get("name", "unknown")
    shape_idx = shape_info.get("shape_index", "?")
    return f"shape_index {shape_idx}('{shape_name}')が{', '.join(issues)}しています"


def describe_shape(shape, index: int) -> dict:
    """1つのshapeの構造情報をdictにする。

    inspect_pptx.py（一覧表示）が使う。ここで定義する `shape_index` は
    `enumerate(slide.shapes)` の0始まり連番であり、edit_pptx.py の各操作が
    指定する `shape_index` と完全に一致する仕様として両スクリプトで揃える
    （edit_pptx.py自体はこの関数を呼ばないが、shape_indexの数え方はここに合わせる）。

    `left_cm`/`top_cm`/`width_cm`/`height_cm` は edit_pptx.py の
    `set_shape_position` が受け取る単位（cm）と完全に一致させている
    （pptx-inspectで読んだ値をそのままset_shape_positionの引数へ使い回せる
    ようにするため。プレースホルダ等でレイアウト側から位置を継承していて
    実座標が取得できない場合はNoneになる）。
    """
    info: dict = {
        "shape_index": index,
        "name": shape.name,
        "shape_type": str(shape.shape_type) if shape.shape_type is not None else None,
        "is_placeholder": shape.is_placeholder,
        "placeholder_idx": None,
        "placeholder_type": None,
        "has_text_frame": shape.has_text_frame,
        "text_preview": None,
        "has_table": shape.has_table,
        "table_dims": None,
        "has_picture": shape.shape_type == MSO_SHAPE_TYPE.PICTURE,
        "left_cm": _length_cm(shape.left),
        "top_cm": _length_cm(shape.top),
        "width_cm": _length_cm(shape.width),
        "height_cm": _length_cm(shape.height),
    }
    if shape.is_placeholder:
        info["placeholder_idx"] = shape.placeholder_format.idx
        info["placeholder_type"] = str(shape.placeholder_format.type)
    if shape.has_text_frame:
        text = shape.text_frame.text
        info["text_preview"] = text[:50] if text else ""
    if shape.has_table:
        table = shape.table
        info["table_dims"] = {"rows": len(table.rows), "cols": len(table.columns)}
    return info


def register_output_path(path, description: str | None = None) -> dict[str, str] | None:
    """生成/更新したファイルをパスメモリーへ登録し、{"@N": 絶対パス} を返す。

    run_script が子プロセスへ注入する AGENT_SRC_DIR 経由で src/path_memory.py
    を import する。AGENT_SRC_DIR未設定やimport失敗時はNone（run_script以外
    から直接実行された場合でもスクリプト自体は失敗させないためのフォールバック）。
    """
    src_dir = os.environ.get("AGENT_SRC_DIR")
    if not src_dir:
        return None
    if src_dir not in sys.path:
        sys.path.insert(0, src_dir)
    try:
        import path_memory
    except ImportError:
        return None
    thread_id, pm_dir, max_entries = path_memory.env_params()
    abs_path = str(Path(path).resolve())
    idx = path_memory.register(thread_id, abs_path, pm_dir, max_entries, description=description)
    if idx is None:
        return None
    return {f"@{idx}": abs_path}


def write_json_result(result: dict, category: str, source_path) -> dict:
    """result全体をJSONファイルへ書き出し、要約に足す {"result_path", "path_memory"} を返す。

    保存先は execute_python_code の中間生成物と同じ `_tmp_<thread_id>/<category>/`
    規約（pdf-tools の render_pdf_pages.py 参照）。会話終了時に自動削除される。
    """
    thread_id = os.environ.get("AGENT_THREAD_ID") or "_no_session"
    out_dir = Path.cwd() / f"_tmp_{thread_id}" / category
    out_dir.mkdir(parents=True, exist_ok=True)
    digest = hashlib.sha1(str(Path(source_path).resolve()).encode("utf-8")).hexdigest()[:8]
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
    out_path = out_dir / f"{digest}_{timestamp}.json"
    out_path.write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")
    info: dict = {"result_path": str(out_path)}
    pm = register_output_path(out_path, description=f"{category}の読み込み結果全文JSON")
    if pm:
        info["path_memory"] = pm
    return info


def summarize_result(result: dict, omit_keys: list[str]) -> dict:
    """result から omit_keys で指定したキーを除いた要約辞書を返す。

    除いたキーがlist型なら "<key>_count"、str型なら "<key>_length" を
    件数/文字数として残す（本文全体はファイル側にのみ残す）。
    """
    summary = dict(result)
    for key in omit_keys:
        if key not in summary:
            continue
        value = summary.pop(key)
        if isinstance(value, list):
            summary[f"{key}_count"] = len(value)
        elif isinstance(value, str):
            summary[f"{key}_length"] = len(value)
    return summary
