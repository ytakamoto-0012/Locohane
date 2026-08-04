"""pptx-tools スキル内の各スクリプトが共有するヘルパー関数。

run_script からは直接実行されない（scripts/ prefix チェックを通らないため）。
scripts/ 配下の各スクリプトが同一ディレクトリから import して使う。
"""

from __future__ import annotations

import os
import shutil
import sys
from datetime import datetime
from pathlib import Path

from pptx.enum.shapes import MSO_SHAPE_TYPE


def setup_utf8_stdio() -> None:
    """標準出力/標準エラーをUTF-8に固定する。

    Windows環境ではパイプ経由のPython子プロセスの既定エンコーディングが
    システムのANSIコードページ（例: cp932）になり、run_script側の
    encoding="utf-8" デコードと食い違って日本語が文字化けするため、
    各スクリプトの先頭で必ず呼ぶこと。
    """
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")


def backup_before_overwrite(path: Path) -> Path | None:
    """path が既に存在する場合、上書き直前に同じフォルダへタイムスタンプ付きで
    コピーしてバックアップを作成する。存在しなければ何もせず None を返す
    （新規作成時はバックアップ不要）。
    """
    if not path.exists():
        return None
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    backup_path = path.with_name(f"{path.stem}.bak_{timestamp}{path.suffix}")
    suffix_n = 2
    while backup_path.exists():
        backup_path = path.with_name(f"{path.stem}.bak_{timestamp}_{suffix_n}{path.suffix}")
        suffix_n += 1
    shutil.copy2(path, backup_path)
    return backup_path


def describe_shape(shape, index: int) -> dict:
    """1つのshapeの構造情報をdictにする。

    inspect_pptx.py（一覧表示）と edit_pptx.py（エラーメッセージ・対象種別
    チェック）の両方から使う共有ロジック。ここで定義する `shape_index` は
    `enumerate(slide.shapes)` の0始まり連番であり、edit_pptx.py の各操作が
    指定する `shape_index` と完全に一致する仕様として両スクリプトで揃える。
    """
    info: dict = {
        "shape_index": index,
        "name": shape.name,
        "shape_type": str(shape.shape_type) if shape.shape_type is not None else None,
        "is_placeholder": shape.is_placeholder,
        "placeholder_idx": None,
        "placeholder_type": None,
        "has_text_frame": shape.has_text_frame,
        "text_preview": None,
        "has_table": shape.has_table,
        "table_dims": None,
        "has_picture": shape.shape_type == MSO_SHAPE_TYPE.PICTURE,
    }
    if shape.is_placeholder:
        info["placeholder_idx"] = shape.placeholder_format.idx
        info["placeholder_type"] = str(shape.placeholder_format.type)
    if shape.has_text_frame:
        text = shape.text_frame.text
        info["text_preview"] = text[:50] if text else ""
    if shape.has_table:
        table = shape.table
        info["table_dims"] = {"rows": len(table.rows), "cols": len(table.columns)}
    return info


def register_output_path(path, description: str | None = None) -> dict[str, str] | None:
    """生成/更新したファイルをパスメモリーへ登録し、{"@N": 絶対パス} を返す。

    run_script が子プロセスへ注入する AGENT_SRC_DIR 経由で src/path_memory.py
    を import する。AGENT_SRC_DIR未設定やimport失敗時はNone（run_script以外
    から直接実行された場合でもスクリプト自体は失敗させないためのフォールバック）。
    """
    src_dir = os.environ.get("AGENT_SRC_DIR")
    if not src_dir:
        return None
    if src_dir not in sys.path:
        sys.path.insert(0, src_dir)
    try:
        import path_memory
    except ImportError:
        return None
    thread_id, pm_dir, max_entries = path_memory.env_params()
    abs_path = str(Path(path).resolve())
    idx = path_memory.register(thread_id, abs_path, pm_dir, max_entries, description=description)
    if idx is None:
        return None
    return {f"@{idx}": abs_path}
