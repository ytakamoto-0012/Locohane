"""既存pptx(PowerPoint)テンプレートに操作列を適用し、別ファイルとして保存する。

pptx-tools スキルの実行スクリプト（progressive disclosure 第3段階）。
run_script ツールから
    python edit_pptx.py <template_path> <output_path> (--data JSON | --data-file PATH) [--overwrite]
の形で呼ばれる。

このプロジェクトのローカルエージェントは任意パスを自由にRead/Editできる
汎用ツールを持たない（run_script はこのスキル配下のスクリプト実行のみ）ため、
Anthropic公式pptxスキルのような「生XMLをエージェントが手で編集する」方式は
取れない。代わりに、構造化されたJSON操作列（テンプレートの何をどう変えるか）を
1回のスクリプト実行で丸ごと適用する方式を取る。テンプレートの意匠
（テーマ・マスター・レイアウト）は変更せず、既存スライドの中身だけを
書き換える／複製・削除・並び替えする。
"""

from __future__ import annotations

import argparse
import copy
import io
import json
import sys
from pathlib import Path

from _common import setup_utf8_stdio

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.exc import PackageNotFoundError
from pptx.oxml.ns import qn

# duplicate_slide で単純なXML deep copyでは正しく複製できないshape種別。
# （チャート/OLE/動画/SmartArt/グループはリレーションシップや埋め込みデータの
# 再登録が必要で、単純コピーだと壊れたpptxになるため明示的にエラーにする。）
UNSUPPORTED_DUPLICATE_TYPES = {
    MSO_SHAPE_TYPE.CHART,
    MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
    MSO_SHAPE_TYPE.LINKED_OLE_OBJECT,
    MSO_SHAPE_TYPE.MEDIA,
    MSO_SHAPE_TYPE.WEB_VIDEO,
    MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT,
    MSO_SHAPE_TYPE.DIAGRAM,
    MSO_SHAPE_TYPE.IGX_GRAPHIC,
    MSO_SHAPE_TYPE.GROUP,
}


def _require(op: dict, key: str):
    value = op.get(key)
    if value is None:
        raise ValueError(f"'{key}' が必要です")
    return value


def _get_slide(prs: Presentation, slide_num: int):
    idx = int(slide_num) - 1
    total = len(prs.slides)
    if idx < 0 or idx >= total:
        raise ValueError(f"存在しないスライド番号です: {slide_num}（現在の総スライド数: {total}）")
    return prs.slides[idx]


def _get_shape(slide, shape_index: int):
    shapes = list(slide.shapes)
    idx = int(shape_index)
    if idx < 0 or idx >= len(shapes):
        raise ValueError(f"存在しないshape_indexです: {shape_index}（このスライドのshape数: {len(shapes)}）")
    return shapes[idx]


def _set_bullets(text_frame, bullets: list) -> None:
    tf = text_frame
    tf.clear()
    first = True
    for item in bullets:
        if isinstance(item, dict):
            text = str(item.get("text", ""))
            level = int(item.get("level", 0))
        else:
            text = str(item)
            level = 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = text
        p.level = level


def op_set_title(prs: Presentation, op: dict) -> None:
    slide = _get_slide(prs, _require(op, "slide"))
    if slide.shapes.title is None:
        raise ValueError("このスライドにタイトルプレースホルダがありません")
    slide.shapes.title.text = str(op.get("text", ""))


def op_set_text(prs: Presentation, op: dict) -> None:
    slide = _get_slide(prs, _require(op, "slide"))
    shape = _get_shape(slide, _require(op, "shape_index"))
    if not shape.has_text_frame:
        raise ValueError(f"shape_index {op['shape_index']} はテキストを持てないshapeです")
    bullets = _require(op, "bullets")
    _set_bullets(shape.text_frame, bullets)


def op_set_table_cell(prs: Presentation, op: dict) -> None:
    slide = _get_slide(prs, _require(op, "slide"))
    shape = _get_shape(slide, _require(op, "shape_index"))
    if not shape.has_table:
        raise ValueError(f"shape_index {op['shape_index']} は表ではありません")
    table = shape.table
    row = int(_require(op, "row"))
    col = int(_require(op, "col"))
    if row < 0 or row >= len(table.rows) or col < 0 or col >= len(table.columns):
        raise ValueError(
            f"表の範囲外です（既存の表は{len(table.rows)}行{len(table.columns)}列、指定は row={row}, col={col}）"
        )
    table.cell(row, col).text = str(op.get("text", ""))


def op_set_table(prs: Presentation, op: dict) -> None:
    slide = _get_slide(prs, _require(op, "slide"))
    shape = _get_shape(slide, _require(op, "shape_index"))
    if not shape.has_table:
        raise ValueError(f"shape_index {op['shape_index']} は表ではありません")
    table = shape.table
    headers = op.get("headers") or []
    rows = op.get("rows") or []
    n_cols_needed = len(headers) if headers else (len(rows[0]) if rows else 0)
    n_rows_needed = (1 if headers else 0) + len(rows)
    if n_rows_needed != len(table.rows) or n_cols_needed != len(table.columns):
        raise ValueError(
            "既存の表の行列数と一致しません"
            f"（既存: {len(table.rows)}行{len(table.columns)}列 / 指定: {n_rows_needed}行{n_cols_needed}列）。"
            "python-pptxは既存表の行列数の増減に対応していないため、"
            "行列数を変えたい場合はcreate_pptx.pyで新規スライドとして作り直してください。"
        )
    r = 0
    if headers:
        for c, h in enumerate(headers):
            table.cell(0, c).text = str(h)
        r = 1
    for row_values in rows:
        for c, val in enumerate(row_values):
            table.cell(r, c).text = str(val)
        r += 1


def op_set_notes(prs: Presentation, op: dict) -> None:
    slide = _get_slide(prs, _require(op, "slide"))
    slide.notes_slide.notes_text_frame.text = str(op.get("text", ""))


def op_replace_picture(prs: Presentation, op: dict) -> None:
    slide = _get_slide(prs, _require(op, "slide"))
    shape = _get_shape(slide, _require(op, "shape_index"))
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        raise ValueError(f"shape_index {op['shape_index']} は画像ではありません")
    image_path = _require(op, "image_path")
    image_file = Path(str(image_path))
    if not image_file.is_file():
        raise ValueError(f"image_path が見つかりません: {image_path}")
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    shape._element.getparent().remove(shape._element)
    slide.shapes.add_picture(str(image_file), left, top, width=width, height=height)


def _copy_picture_into(new_slide, shape) -> None:
    image = shape.image
    _, rid = new_slide.part.get_or_add_image_part(io.BytesIO(image.blob))
    new_el = copy.deepcopy(shape._element)
    blip = new_el.find(f".//{qn('a:blip')}")
    if blip is not None:
        blip.set(qn("r:embed"), rid)
    new_slide.shapes._spTree.append(new_el)


def _duplicate_slide_once(prs: Presentation, source_slide):
    layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)
    # add_slide が layout に合わせて自動生成する空プレースホルダは、
    # source_slide の内容をそのままコピーし直すため一旦すべて取り除く。
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)

    for shape in source_slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            _copy_picture_into(new_slide, shape)
        elif shape.shape_type in UNSUPPORTED_DUPLICATE_TYPES:
            raise ValueError(
                f"複製元スライドに非対応のshape種別（{shape.shape_type}, name={shape.name!r}）"
                "が含まれるため複製できません（チャート/OLE/動画/SmartArt/グループは非対応）"
            )
        else:
            # プレースホルダ・表・テキストボックス等はXML要素をそのままdeep copyすれば
            # 書式・placeholder紐付け（layout側のidx/type一致）ともに維持される。
            new_slide.shapes._spTree.append(copy.deepcopy(shape._element))

    if source_slide.has_notes_slide:
        notes_text = source_slide.notes_slide.notes_text_frame.text
        if notes_text.strip():
            new_slide.notes_slide.notes_text_frame.text = notes_text

    return new_slide


def _move_slide_to(prs: Presentation, position: int) -> None:
    """add_slide で末尾に追加された直後のスライドを0始まりposition位置へ移動する。"""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    last = slides[-1]
    xml_slides.remove(last)
    xml_slides.insert(position, last)


def op_duplicate_slide(prs: Presentation, op: dict) -> None:
    source_num = int(_require(op, "slide"))
    source_slide = _get_slide(prs, source_num)
    insert_after = int(op.get("insert_after", source_num))
    count = int(op.get("count", 1))
    if count < 1:
        raise ValueError("count は1以上を指定してください")
    total = len(prs.slides)
    if insert_after < 0 or insert_after > total:
        raise ValueError(f"insert_after が不正です（0〜{total}の範囲で指定してください）: {insert_after}")

    position = insert_after
    for _ in range(count):
        _duplicate_slide_once(prs, source_slide)
        _move_slide_to(prs, position)
        position += 1


def op_delete_slide(prs: Presentation, op: dict) -> None:
    slide_num = int(_require(op, "slide"))
    idx = slide_num - 1
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    if idx < 0 or idx >= len(slides):
        raise ValueError(f"存在しないスライド番号です: {slide_num}（現在の総スライド数: {len(slides)}）")
    slide_id_elm = slides[idx]
    rid = slide_id_elm.get(qn("r:id"))
    prs.part.drop_rel(rid)
    xml_slides.remove(slide_id_elm)


def op_reorder_slides(prs: Presentation, op: dict) -> None:
    order = op.get("order")
    total = len(prs.slides)
    if not isinstance(order, list) or sorted(order) != list(range(1, total + 1)):
        raise ValueError(
            f"order には現在の全スライド番号(1〜{total})の順列を指定してください: {order!r}"
        )
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    new_order_elms = [slides[i - 1] for i in order]
    for elm in new_order_elms:
        xml_slides.remove(elm)
    for elm in new_order_elms:
        xml_slides.append(elm)


OP_HANDLERS = {
    "set_title": op_set_title,
    "set_text": op_set_text,
    "set_table_cell": op_set_table_cell,
    "set_table": op_set_table,
    "set_notes": op_set_notes,
    "replace_picture": op_replace_picture,
    "duplicate_slide": op_duplicate_slide,
    "delete_slide": op_delete_slide,
    "reorder_slides": op_reorder_slides,
}


def main() -> int:
    setup_utf8_stdio()
    parser = argparse.ArgumentParser()
    parser.add_argument("template_path")
    parser.add_argument("output_path")
    group = parser.add_mutually_exclusive_group(required=True)
    group.add_argument("--data")
    group.add_argument("--data-file")
    parser.add_argument("--overwrite", action="store_true")
    args = parser.parse_args()

    template_path = Path(args.template_path)
    output_path = Path(args.output_path)

    if not template_path.exists():
        print(f"テンプレートファイルが見つかりません: {args.template_path}", file=sys.stderr)
        return 1
    if template_path.is_dir():
        print(f"指定パスはディレクトリです（ファイル専用）: {args.template_path}", file=sys.stderr)
        return 1
    if template_path.resolve() == output_path.resolve() and not args.overwrite:
        print(
            "テンプレートと同じパスへの保存は --overwrite を指定した場合のみ許可されます。"
            "別のoutput_pathを指定するか、上書きしてよいことをユーザーに確認してから --overwrite を付けてください。",
            file=sys.stderr,
        )
        return 1

    if args.data_file:
        data_file = Path(args.data_file)
        if not data_file.exists():
            print(f"データファイルが見つかりません: {args.data_file}", file=sys.stderr)
            return 1
        try:
            raw = data_file.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            raw = data_file.read_text(encoding="cp932")
    else:
        raw = args.data

    try:
        data = json.loads(raw)
    except json.JSONDecodeError as e:
        print(f"JSONとして解析できませんでした: {e}", file=sys.stderr)
        return 1

    operations = data.get("operations") if isinstance(data, dict) else None
    if not operations:
        print("data の 'operations' キーに1件以上の操作定義が必要です。", file=sys.stderr)
        return 1

    try:
        prs = Presentation(str(template_path))
    except PackageNotFoundError:
        print(f"pptxとして読み込めませんでした（壊れているか非対応形式の可能性）: {args.template_path}", file=sys.stderr)
        return 1
    except Exception as e:  # noqa: BLE001 - python-pptx側の想定外エラーもエラー終了に統一する
        print(f"pptxの読み込みに失敗しました: {e}", file=sys.stderr)
        return 1

    for i, op in enumerate(operations):
        op_name = op.get("op") if isinstance(op, dict) else None
        handler = OP_HANDLERS.get(op_name)
        if handler is None:
            supported = ", ".join(sorted(OP_HANDLERS))
            print(f"操作{i + 1}: 未対応の op です: {op_name!r}（対応: {supported}）", file=sys.stderr)
            return 1
        try:
            handler(prs, op)
        except ValueError as e:
            print(f"操作{i + 1}（{op_name}）の適用に失敗しました: {e}", file=sys.stderr)
            return 1

    output_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(str(output_path))
    except OSError as e:
        print(f"ファイルの保存に失敗しました: {e}", file=sys.stderr)
        return 1

    result = {
        "output_path": str(output_path),
        "total_slides": len(prs.slides),
        "size_bytes": output_path.stat().st_size,
        "applied_operations": len(operations),
    }
    print(json.dumps(result, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    sys.exit(main())
