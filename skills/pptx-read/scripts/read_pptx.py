"""pptx(PowerPoint)ファイルからテキスト・表・ノートを抽出して JSON で出力する。

pptx-tools スキルの実行スクリプト（progressive disclosure 第3段階）。
run_script ツールから
    python read_pptx.py <pptx_path> [--start-slide N] [--max-slides N]
の形で呼ばれる。

read_pdf.py の start-page/max-pages と同じページング設計をスライド単位で適用し、
1回の呼び出しで返すスライド数の上限を設けて巨大な出力を避ける。
"""

from __future__ import annotations

import argparse
import json
import sys
import traceback
from pathlib import Path

# office_shared/pptx_common.py から共有ヘルパーを import する（1-B 相互import方式）
_OFFICE_SHARED = Path(__file__).resolve().parent.parent.parent / "office_shared"
if str(_OFFICE_SHARED) not in sys.path:
    sys.path.append(str(_OFFICE_SHARED))
from pptx_common import setup_utf8_stdio, summarize_result, write_json_result  # noqa: E402

from pptx import Presentation
from pptx.exc import PackageNotFoundError


def extract_slide(slide) -> dict:
    title_shape = slide.shapes.title
    title = title_shape.text_frame.text if title_shape is not None and title_shape.has_text_frame else None
    # slide.shapes.title と slide.shapes の反復で得られるオブジェクトは、python-pptx内部で
    # 都度生成されるラッパーのため is 比較では同一視できない。shape_id で突き合わせる。
    title_id = title_shape.shape_id if title_shape is not None else None

    texts: list[str] = []
    tables: list[list[list[str]]] = []
    for shape in slide.shapes:
        if title_id is not None and shape.shape_id == title_id:
            continue
        if shape.has_text_frame:
            text = shape.text_frame.text
            if text.strip():
                texts.append(text)
        if shape.has_table:
            table_data = [[cell.text for cell in row.cells] for row in shape.table.rows]
            tables.append(table_data)

    notes = None
    if slide.has_notes_slide:
        notes_text = slide.notes_slide.notes_text_frame.text
        notes = notes_text if notes_text.strip() else None

    return {"title": title, "texts": texts, "tables": tables, "notes": notes}


def main() -> int:
    setup_utf8_stdio()
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx_path")
    parser.add_argument("--start-slide", type=int, default=1)
    parser.add_argument("--max-slides", type=int, default=20)
    args = parser.parse_args()

    path = Path(args.pptx_path)
    if not path.exists():
        print(f"ファイルが見つかりません: {args.pptx_path}", file=sys.stderr)
        return 1
    if path.is_dir():
        print(f"指定パスはディレクトリです（ファイル専用）: {args.pptx_path}", file=sys.stderr)
        return 1

    try:
        prs = Presentation(str(path))
    except PackageNotFoundError:
        print(f"pptxとして読み込めませんでした（壊れているか非対応形式の可能性）: {args.pptx_path}", file=sys.stderr)
        return 1
    except Exception as e:  # noqa: BLE001 - python-pptx側の想定外エラーもエラー終了に統一する
        print(f"pptxの読み込みに失敗しました: {e}\n{traceback.format_exc()}", file=sys.stderr)
        return 1

    total_slides = len(prs.slides)

    start_idx = max(args.start_slide, 1) - 1
    end_idx = min(start_idx + max(args.max_slides, 1), total_slides)

    slides_out = []
    for i in range(start_idx, end_idx):
        info = extract_slide(prs.slides[i])
        slides_out.append({"index": i + 1, **info})

    result = {
        "path": str(path),
        "total_slides": total_slides,
        "start_slide": start_idx + 1 if slides_out else None,
        "end_slide": end_idx if slides_out else None,
        "slides": slides_out,
    }
    summary = summarize_result(result, ["slides"])
    summary.update(write_json_result(result, "pptx_read", path))
    print(json.dumps(summary, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    sys.exit(main())
