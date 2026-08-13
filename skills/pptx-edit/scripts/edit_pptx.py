"""既存pptx(PowerPoint)テンプレートに操作列を適用し、別ファイルとして保存する。

pptx-tools スキルの実行スクリプト（progressive disclosure 第3段階）。
run_script ツールから
    python edit_pptx.py <template_path> <output_path> (--data JSON | --data-file PATH) [--overwrite]
の形で呼ばれる。

このプロジェクトのローカルエージェントは任意パスを自由にRead/Editできる
汎用ツールを持たない（run_script はこのスキル配下のスクリプト実行のみ）ため、
Anthropic公式pptxスキルのような「生XMLをエージェントが手で編集する」方式は
取れない。代わりに、構造化されたJSON操作列（テンプレートの何をどう変えるか）を
1回のスクリプト実行で丸ごと適用する方式を取る。テンプレートの意匠
（テーマ・マスター・レイアウト）は変更せず、既存スライドの中身だけを
書き換える／複製・削除・並び替えする。
"""

from __future__ import annotations

import argparse
import copy
import io
import json
import sys
import traceback
from pathlib import Path

from _common import backup_before_overwrite, register_output_path, setup_utf8_stdio

# _shared/office_theme.py から THEMES / resolve_theme を import する（1-B 相互import方式）
_OFFICE_SHARED = Path(__file__).resolve().parent.parent.parent / "_shared"
if str(_OFFICE_SHARED) not in sys.path:
    sys.path.append(str(_OFFICE_SHARED))
from office_theme import THEMES, resolve_theme  # noqa: E402
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.exc import PackageNotFoundError
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Cm, Pt

# duplicate_slide で単純なXML deep copyでは正しく複製できないshape種別。
# （チャート/OLE/動画/SmartArt/グループはリレーションシップや埋め込みデータの
# 再登録が必要で、単純コピーだと壊れたpptxになるため明示的にエラーにする。）
UNSUPPORTED_DUPLICATE_TYPES = {
    MSO_SHAPE_TYPE.CHART,
    MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
    MSO_SHAPE_TYPE.LINKED_OLE_OBJECT,
    MSO_SHAPE_TYPE.MEDIA,
    MSO_SHAPE_TYPE.WEB_VIDEO,
    MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT,
    MSO_SHAPE_TYPE.DIAGRAM,
    MSO_SHAPE_TYPE.IGX_GRAPHIC,
    MSO_SHAPE_TYPE.GROUP,
}


def _require(op: dict, key: str):
    value = op.get(key)
    if value is None:
        raise ValueError(f"'{key}' が必要です")
    return value


class EditContext:
    """opsバッチ全体で共有する状態。

    `slide`番号は常に「このopsバッチを開始した時点（＝直前のpptx-inspectが
    見せていたはずの状態）のスライド番号」として解決する。`duplicate_slide`/
    `delete_slide`/`reorder_slides`が後続スライド番号をずらしても、それより
    後のopで毎回シフト量を手計算する必要が無い（Excel編集での同種の事故
    ＝行挿入後の絶対行番号ズレ、と同根の問題への対応）。

    対して`duplicate_slide`の`insert_after`と`reorder_slides`の`order`は、
    「新しい内容をどこに置くか／並べ替えるか」という現在時点の配置に関する
    指定なので、現在のライブ位置基準のまま扱う（詳細はSKILL.md参照）。
    """

    def __init__(self, prs: Presentation):
        self.prs = prs
        # sldId要素はスライドの並び替え・複製では再生成されず同一オブジェクトの
        # ままなので、Python識別性でのlist.index()による生存追跡が可能。
        # 削除された要素だけがこのリストの走査から消える。
        self._original_slide_elements = list(prs.slides._sldIdLst)

    def _resolve_element(self, slide_num: int):
        idx = int(slide_num) - 1
        total = len(self._original_slide_elements)
        if idx < 0 or idx >= total:
            raise ValueError(f"存在しないスライド番号です: {slide_num}（このopsバッチ開始時点の総スライド数: {total}）")
        element = self._original_slide_elements[idx]
        if element not in list(self.prs.slides._sldIdLst):
            raise ValueError(f"スライド{slide_num}は、この呼び出し内の先行する操作（delete_slide等）で" "既に削除されています")
        return element

    def get_slide(self, slide_num: int):
        element = self._resolve_element(slide_num)
        live_idx = list(self.prs.slides._sldIdLst).index(element)
        return self.prs.slides[live_idx]

    def get_slide_element(self, slide_num: int):
        return self._resolve_element(slide_num)

    def current_position(self, slide_num: int) -> int:
        """スライド番号（バッチ開始時点基準）を、現在のライブな1始まり位置へ変換する。"""
        element = self._resolve_element(slide_num)
        return list(self.prs.slides._sldIdLst).index(element) + 1


def _get_shape(slide, shape_index: int):
    shapes = list(slide.shapes)
    idx = int(shape_index)
    if idx < 0 or idx >= len(shapes):
        raise ValueError(f"存在しないshape_indexです: {shape_index}（このスライドのshape数: {len(shapes)}）")
    return shapes[idx]


def _set_bullets(text_frame, bullets: list) -> None:
    tf = text_frame
    tf.clear()
    first = True
    for item in bullets:
        if isinstance(item, dict):
            text = str(item.get("text", ""))
            level = int(item.get("level", 0))
        else:
            text = str(item)
            level = 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = text
        p.level = level


def op_set_title(ctx: EditContext, op: dict) -> None:
    slide = ctx.get_slide(_require(op, "slide"))
    if slide.shapes.title is None:
        raise ValueError("このスライドにタイトルプレースホルダがありません")
    slide.shapes.title.text = str(op.get("text", ""))


def op_set_text(ctx: EditContext, op: dict) -> None:
    slide = ctx.get_slide(_require(op, "slide"))
    shape = _get_shape(slide, _require(op, "shape_index"))
    if not shape.has_text_frame:
        raise ValueError(f"shape_index {op['shape_index']} はテキストを持てないshapeです")
    bullets = _require(op, "bullets")
    _set_bullets(shape.text_frame, bullets)


def op_set_table_cell(ctx: EditContext, op: dict) -> None:
    slide = ctx.get_slide(_require(op, "slide"))
    shape = _get_shape(slide, _require(op, "shape_index"))
    if not shape.has_table:
        raise ValueError(f"shape_index {op['shape_index']} は表ではありません")
    table = shape.table
    row = int(_require(op, "row"))
    col = int(_require(op, "col"))
    if row < 0 or row >= len(table.rows) or col < 0 or col >= len(table.columns):
        raise ValueError(f"表の範囲外です（既存の表は{len(table.rows)}行{len(table.columns)}列、指定は row={row}, col={col}）")
    table.cell(row, col).text = str(op.get("text", ""))


def op_set_table(ctx: EditContext, op: dict) -> None:
    slide = ctx.get_slide(_require(op, "slide"))
    shape = _get_shape(slide, _require(op, "shape_index"))
    if not shape.has_table:
        raise ValueError(f"shape_index {op['shape_index']} は表ではありません")
    table = shape.table
    headers = op.get("headers") or []
    rows = op.get("rows") or []
    n_cols_needed = len(headers) if headers else (len(rows[0]) if rows else 0)
    n_rows_needed = (1 if headers else 0) + len(rows)
    if n_rows_needed != len(table.rows) or n_cols_needed != len(table.columns):
        raise ValueError(
            "既存の表の行列数と一致しません"
            f"（既存: {len(table.rows)}行{len(table.columns)}列 / 指定: {n_rows_needed}行{n_cols_needed}列）。"
            "python-pptxは既存表の行列数の増減に対応していないため、"
            "行列数を変えたい場合はcreate_pptx.pyで新規スライドとして作り直してください。"
        )
    r = 0
    if headers:
        for c, h in enumerate(headers):
            table.cell(0, c).text = str(h)
        r = 1
    for row_values in rows:
        for c, val in enumerate(row_values):
            table.cell(r, c).text = str(val)
        r += 1


def op_set_notes(ctx: EditContext, op: dict) -> None:
    slide = ctx.get_slide(_require(op, "slide"))
    slide.notes_slide.notes_text_frame.text = str(op.get("text", ""))


def op_replace_picture(ctx: EditContext, op: dict) -> None:
    slide = ctx.get_slide(_require(op, "slide"))
    shape = _get_shape(slide, _require(op, "shape_index"))
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        raise ValueError(f"shape_index {op['shape_index']} は画像ではありません")
    image_path = _require(op, "image_path")
    image_file = Path(str(image_path))
    if not image_file.is_file():
        raise ValueError(f"image_path が見つかりません: {image_path}")
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    shape._element.getparent().remove(shape._element)
    slide.shapes.add_picture(str(image_file), left, top, width=width, height=height)


_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}


def _set_run_font_name(run, font_name: str) -> None:
    # run.font.name はlatin書体のみを設定するため、日本語等の実際の描画に
    # 使われる east-asian書体（a:ea）もoxmlで直接設定する（docx-editの
    # _set_east_asian_font と同じ理由）。
    run.font.name = font_name
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = OxmlElement("a:ea")
        rPr.append(ea)
    ea.set("typeface", font_name)


def _style_text_frame(
    text_frame,
    text_color: str | None,
    bold: bool | None,
    font_size_pt,
    italic: bool | None = None,
    underline: bool | None = None,
    font_name: str | None = None,
    align: str | None = None,
) -> None:
    for paragraph in text_frame.paragraphs:
        if align:
            paragraph.alignment = _ALIGN_MAP[align]
        for run in paragraph.runs:
            if text_color:
                run.font.color.rgb = RGBColor.from_string(text_color)
            if bold is not None:
                run.font.bold = bold
            if italic is not None:
                run.font.italic = italic
            if underline is not None:
                run.font.underline = underline
            if font_size_pt:
                run.font.size = Pt(float(font_size_pt))
            if font_name:
                _set_run_font_name(run, font_name)


def op_set_shape_style(ctx: EditContext, op: dict) -> None:
    """既存shape（テキストボックス・オートシェイプ・表）を明示的に再配色・再配置する。

    ユーザーが「もっと格好よく／見やすくして」「文字を中央寄せに」等、既存
    デザインの変更を明示的に頼んできた場合にのみ使う（set_title等の中身
    差し替えopは見た目に一切触れない設計なので、再配色・配置変更が必要な
    場面はこのopに限定する）。図形自体の位置/サイズ移動は set_shape_position
    を使う。
    """
    slide = ctx.get_slide(_require(op, "slide"))
    shape = _get_shape(slide, _require(op, "shape_index"))

    theme = resolve_theme(op["theme"]) if op.get("theme") else None
    role = op.get("role")
    text_color = op.get("text_color")
    bold = op.get("bold")
    italic = op.get("italic")
    underline = op.get("underline")
    font_size_pt = op.get("font_size_pt")
    font_name = op.get("font_name")
    fill_color = op.get("fill_color")
    border_color = op.get("border_color")
    align = op.get("align")

    if align is not None and align not in _ALIGN_MAP:
        raise ValueError(f"未対応の align です: {align!r}（対応: {', '.join(_ALIGN_MAP)}）")

    if role == "heading":
        text_color = text_color or (theme["primary"] if theme else None)
        bold = True if bold is None else bold
    elif role == "table_header":
        fill_color = fill_color or (theme["primary"] if theme else None)
        text_color = text_color or (theme["text_on_primary"] if theme else None)
        bold = True if bold is None else bold
    elif role is not None:
        raise ValueError(f"未対応の role です: {role!r}（対応: heading, table_header）")

    if not any([text_color, bold is not None, italic is not None, underline is not None, font_size_pt, font_name, fill_color, border_color, align]):
        raise ValueError(
            "text_color/bold/italic/underline/font_size_pt/font_name/fill_color/"
            "border_color/align のいずれか、または role（heading/table_header）+ theme の指定が必要です"
        )

    if shape.has_table:
        if border_color:
            raise ValueError("表shapeの罫線色（border_color）は現状非対応です（セル塗り・文字装飾のみ対応）")
        table = shape.table
        if role == "table_header":
            target_rows = [0]
        elif op.get("all_rows"):
            target_rows = range(len(table.rows))
        else:
            row = op.get("row")
            if row is None:
                raise ValueError(
                    "表shapeへの set_shape_style は対象行の指定が必要です"
                    "（role='table_header' で見出し行、'row'で行番号指定、'all_rows':trueで全行）"
                )
            target_rows = [int(row)]
        for r in target_rows:
            if r < 0 or r >= len(table.rows):
                raise ValueError(f"存在しない行です: {r}（この表は{len(table.rows)}行）")
            for c in range(len(table.columns)):
                cell = table.cell(r, c)
                if fill_color:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor.from_string(fill_color)
                _style_text_frame(
                    cell.text_frame, text_color, bold, font_size_pt, italic=italic, underline=underline, font_name=font_name, align=align
                )
    else:
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(fill_color)
        if border_color:
            shape.line.color.rgb = RGBColor.from_string(border_color)
        if any([text_color, bold is not None, italic is not None, underline is not None, font_size_pt, font_name, align]):
            if not shape.has_text_frame:
                raise ValueError(f"shape_index {op['shape_index']} はテキストを持てないshapeです")
            _style_text_frame(shape.text_frame, text_color, bold, font_size_pt, italic=italic, underline=underline, font_name=font_name, align=align)


def op_set_shape_position(ctx: EditContext, op: dict) -> None:
    """既存shapeの位置・サイズを変更する（cm単位）。

    ユーザーが「もっと左に」「大きくして」等、配置の変更を明示的に頼んできた
    場合にのみ使う。left/top/width/height のうち指定されたものだけを変更し、
    省略したものは元の値を維持する。
    """
    slide = ctx.get_slide(_require(op, "slide"))
    shape = _get_shape(slide, _require(op, "shape_index"))

    left = op.get("left_cm")
    top = op.get("top_cm")
    width = op.get("width_cm")
    height = op.get("height_cm")

    if left is None and top is None and width is None and height is None:
        raise ValueError("left_cm/top_cm/width_cm/height_cm のいずれか1つ以上の指定が必要です")

    if left is not None:
        shape.left = Cm(float(left))
    if top is not None:
        shape.top = Cm(float(top))
    if width is not None:
        shape.width = Cm(float(width))
    if height is not None:
        shape.height = Cm(float(height))


def _copy_picture_into(new_slide, shape) -> None:
    image = shape.image
    _, rid = new_slide.part.get_or_add_image_part(io.BytesIO(image.blob))
    new_el = copy.deepcopy(shape._element)
    blip = new_el.find(f".//{qn('a:blip')}")
    if blip is not None:
        blip.set(qn("r:embed"), rid)
    new_slide.shapes._spTree.append(new_el)


def _duplicate_slide_once(prs: Presentation, source_slide):
    layout = source_slide.slide_layout
    # python-pptxのPresentationPart.add_slideは新規パート名を
    # 「/ppt/slides/slide<現在のスライド数+1>.xml」という単純な個数ベースで
    # 決めており、実際に空いているか（package.next_partnameのような
    # 網羅チェック）を確認しない。そのため delete_slide で歯抜けができた後に
    # duplicate_slide すると、末尾より手前に残っている既存スライドと
    # パート名が衝突し、保存時にZIP内で同名エントリが重複する
    # （python-pptx側の既知の制限）。add_slide前に衝突しない正しい空き
    # パート名を計算しておき、add_slide後に付け替えて回避する。
    safe_partname = prs.part.package.next_partname("/ppt/slides/slide%d.xml")
    new_slide = prs.slides.add_slide(layout)
    if new_slide.part.partname != safe_partname:
        new_slide.part.partname = safe_partname
    # add_slide が layout に合わせて自動生成する空プレースホルダは、
    # source_slide の内容をそのままコピーし直すため一旦すべて取り除く。
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)

    for shape in source_slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            _copy_picture_into(new_slide, shape)
        elif shape.shape_type in UNSUPPORTED_DUPLICATE_TYPES:
            raise ValueError(
                f"複製元スライドに非対応のshape種別（{shape.shape_type}, name={shape.name!r}）"
                "が含まれるため複製できません（チャート/OLE/動画/SmartArt/グループは非対応）"
            )
        else:
            # プレースホルダ・表・テキストボックス等はXML要素をそのままdeep copyすれば
            # 書式・placeholder紐付け（layout側のidx/type一致）ともに維持される。
            new_slide.shapes._spTree.append(copy.deepcopy(shape._element))

    if source_slide.has_notes_slide:
        notes_text = source_slide.notes_slide.notes_text_frame.text
        if notes_text.strip():
            new_slide.notes_slide.notes_text_frame.text = notes_text

    return new_slide


def _move_slide_to(prs: Presentation, position: int) -> None:
    """add_slide で末尾に追加された直後のスライドを0始まりposition位置へ移動する。"""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    last = slides[-1]
    xml_slides.remove(last)
    xml_slides.insert(position, last)


def op_duplicate_slide(ctx: EditContext, op: dict) -> None:
    source_num = int(_require(op, "slide"))
    source_slide = ctx.get_slide(source_num)
    # insert_after省略時は「複製元スライドの“現在の”ライブ位置の直後」に置く。
    # source_numをそのまま使うと、このバッチ内で先行opによりスライドが
    # 既にずれている場合に誤った位置へ挿入してしまう。
    insert_after = int(op["insert_after"]) if "insert_after" in op else ctx.current_position(source_num)
    count = int(op.get("count", 1))
    if count < 1:
        raise ValueError("count は1以上を指定してください")
    total = len(ctx.prs.slides)
    if insert_after < 0 or insert_after > total:
        raise ValueError(f"insert_after が不正です（0〜{total}の範囲で指定してください）: {insert_after}")

    position = insert_after
    for _ in range(count):
        _duplicate_slide_once(ctx.prs, source_slide)
        _move_slide_to(ctx.prs, position)
        position += 1


def op_delete_slide(ctx: EditContext, op: dict) -> None:
    slide_num = int(_require(op, "slide"))
    element = ctx.get_slide_element(slide_num)
    xml_slides = ctx.prs.slides._sldIdLst
    rid = element.get(qn("r:id"))
    ctx.prs.part.drop_rel(rid)
    xml_slides.remove(element)


def op_reorder_slides(ctx: EditContext, op: dict) -> None:
    """全スライドの並び順を指定順に変更する。

    `order`は他のopと異なり「バッチ開始時点の番号」ではなく、この操作を
    適用する時点でのライブなスライド番号(1始まり)の順列を指定する
    （このopの目的自体が「今存在する全スライドをどう並べるか」であり、
    かつ先行opでduplicate_slideが作った新規スライドにはバッチ開始時点の
    番号が存在しないため）。delete_slide/duplicate_slideと同一バッチで
    併用する場合は、reorder_slidesをそれらより後・かつ操作列の最後に
    置くことを推奨する（詳細はSKILL.md参照）。
    """
    order = op.get("order")
    total = len(ctx.prs.slides)
    if not isinstance(order, list) or sorted(order) != list(range(1, total + 1)):
        raise ValueError(f"order には現在の全スライド番号(1〜{total})の順列を指定してください: {order!r}")
    xml_slides = ctx.prs.slides._sldIdLst
    slides = list(xml_slides)
    new_order_elms = [slides[i - 1] for i in order]
    for elm in new_order_elms:
        xml_slides.remove(elm)
    for elm in new_order_elms:
        xml_slides.append(elm)


OP_HANDLERS = {
    "set_title": op_set_title,
    "set_text": op_set_text,
    "set_table_cell": op_set_table_cell,
    "set_table": op_set_table,
    "set_notes": op_set_notes,
    "replace_picture": op_replace_picture,
    "set_shape_style": op_set_shape_style,
    "set_shape_position": op_set_shape_position,
    "duplicate_slide": op_duplicate_slide,
    "delete_slide": op_delete_slide,
    "reorder_slides": op_reorder_slides,
}


def main() -> int:
    setup_utf8_stdio()
    parser = argparse.ArgumentParser()
    parser.add_argument("template_path")
    parser.add_argument("output_path")
    group = parser.add_mutually_exclusive_group(required=True)
    group.add_argument("--data")
    group.add_argument("--data-file")
    parser.add_argument("--overwrite", action="store_true")
    args = parser.parse_args()

    template_path = Path(args.template_path)
    output_path = Path(args.output_path)

    if not template_path.exists():
        print(f"テンプレートファイルが見つかりません: {args.template_path}", file=sys.stderr)
        return 1
    if template_path.is_dir():
        print(f"指定パスはディレクトリです（ファイル専用）: {args.template_path}", file=sys.stderr)
        return 1
    if template_path.resolve() == output_path.resolve() and not args.overwrite:
        print(
            "テンプレートと同じパスへの保存は --overwrite を指定した場合のみ許可されます。"
            "別のoutput_pathを指定するか、上書きしてよいことをユーザーに確認してから --overwrite を付けてください。",
            file=sys.stderr,
        )
        return 1

    if args.data_file:
        data_file = Path(args.data_file)
        if not data_file.exists():
            print(f"データファイルが見つかりません: {args.data_file}", file=sys.stderr)
            return 1
        try:
            raw = data_file.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            raw = data_file.read_text(encoding="cp932")
    else:
        raw = args.data

    try:
        data = json.loads(raw)
    except json.JSONDecodeError as e:
        print(f"JSONとして解析できませんでした: {e}", file=sys.stderr)
        return 1

    operations = data.get("operations") if isinstance(data, dict) else None
    if not operations:
        print("data の 'operations' キーに1件以上の操作定義が必要です。", file=sys.stderr)
        return 1

    try:
        prs = Presentation(str(template_path))
    except PackageNotFoundError:
        print(f"pptxとして読み込めませんでした（壊れているか非対応形式の可能性）: {args.template_path}", file=sys.stderr)
        return 1
    except Exception as e:  # noqa: BLE001 - python-pptx側の想定外エラーもエラー終了に統一する
        print(f"pptxの読み込みに失敗しました: {e}\n{traceback.format_exc()}", file=sys.stderr)
        return 1

    ctx = EditContext(prs)
    for i, op in enumerate(operations):
        op_name = op.get("op") if isinstance(op, dict) else None
        handler = OP_HANDLERS.get(op_name)
        if handler is None:
            supported = ", ".join(sorted(OP_HANDLERS))
            print(f"操作{i + 1}: 未対応の op です: {op_name!r}（対応: {supported}）", file=sys.stderr)
            return 1
        try:
            handler(ctx, op)
        except ValueError as e:
            print(f"操作{i + 1}（{op_name}）の適用に失敗しました: {e}", file=sys.stderr)
            return 1

    output_path.parent.mkdir(parents=True, exist_ok=True)
    backup_path = backup_before_overwrite(output_path)
    try:
        prs.save(str(output_path))
    except OSError as e:
        print(f"ファイルの保存に失敗しました: {e}", file=sys.stderr)
        return 1

    result = {
        "output_path": str(output_path),
        "backup_path": str(backup_path) if backup_path else None,
        "total_slides": len(prs.slides),
        "size_bytes": output_path.stat().st_size,
        "applied_operations": len(operations),
    }
    path_memory = register_output_path(output_path, description="edit_pptxが生成/更新したPPTX")
    if path_memory:
        result["path_memory"] = path_memory
    print(json.dumps(result, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    sys.exit(main())
