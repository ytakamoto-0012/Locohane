"""JSON定義から新規pptx(PowerPoint)ファイルを生成する。

pptx-tools スキルの実行スクリプト（progressive disclosure 第3段階）。
run_script ツールから
    python create_pptx.py <output_path> (--data JSON | --data-file PATH)
の形で呼ばれる。

assets/template_16x9.pptx（16:9・python-pptx既定テーマ準拠）をテンプレートとして使い、
JSON で渡されたスライド定義を順に slide_layouts へ適用していく。
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from _common import setup_utf8_stdio

from pptx import Presentation
from pptx.util import Inches

TEMPLATE_PATH = Path(__file__).parent.parent / "assets" / "template_16x9.pptx"

# layout名 -> テンプレートの slide_layouts インデックス（既定テーマの構成を継承）。
LAYOUT_MAP = {
    "title": 0,
    "content": 1,
    "section": 2,
    "two_content": 3,
    "table": 1,
    "picture": 1,
    "blank": 6,
}


def set_bullets(placeholder, bullets: list) -> None:
    tf = placeholder.text_frame
    tf.clear()
    first = True
    for item in bullets:
        if isinstance(item, dict):
            text = str(item.get("text", ""))
            level = int(item.get("level", 0))
        else:
            text = str(item)
            level = 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = text
        p.level = level


def add_table(slide, table_def: dict, left, top, width, height) -> None:
    headers = table_def.get("headers") or []
    rows = table_def.get("rows") or []
    n_cols = len(headers) if headers else (len(rows[0]) if rows else 0)
    n_rows = (1 if headers else 0) + len(rows)
    if n_cols == 0 or n_rows == 0:
        raise ValueError("table には headers または rows のいずれかが必要です")

    graphic_frame = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = graphic_frame.table
    r = 0
    if headers:
        for c, h in enumerate(headers):
            table.cell(0, c).text = str(h)
        r = 1
    for row in rows:
        for c, val in enumerate(row):
            table.cell(r, c).text = str(val)
        r += 1


def apply_slide(prs: Presentation, slide_def: dict) -> None:
    layout_key = slide_def.get("layout", "content")
    if layout_key not in LAYOUT_MAP:
        supported = ", ".join(sorted(LAYOUT_MAP))
        raise ValueError(f"未対応の layout です: {layout_key!r}（対応: {supported}）")

    layout_idx = LAYOUT_MAP[layout_key]
    try:
        layout = prs.slide_layouts[layout_idx]
    except IndexError as e:
        raise ValueError(f"テンプレートに slide_layouts[{layout_idx}] が存在しません") from e
    slide = prs.slides.add_slide(layout)

    title = slide_def.get("title")
    if title is not None and slide.shapes.title is not None:
        slide.shapes.title.text = str(title)

    if layout_key == "title":
        subtitle = slide_def.get("subtitle")
        if subtitle is not None:
            try:
                slide.placeholders[1].text = str(subtitle)
            except KeyError:
                pass

    elif layout_key == "content":
        bullets = slide_def.get("bullets") or []
        if bullets:
            try:
                set_bullets(slide.placeholders[1], bullets)
            except KeyError as e:
                raise ValueError("content レイアウトに本文プレースホルダが見つかりません") from e

    elif layout_key == "two_content":
        left_bullets = slide_def.get("left_bullets") or []
        right_bullets = slide_def.get("right_bullets") or []
        try:
            if left_bullets:
                set_bullets(slide.placeholders[1], left_bullets)
            if right_bullets:
                set_bullets(slide.placeholders[2], right_bullets)
        except KeyError as e:
            raise ValueError("two_content レイアウトに左右のプレースホルダが見つかりません") from e

    elif layout_key == "table":
        table_def = slide_def.get("table")
        if not table_def:
            raise ValueError("layout が table のスライドには table キーが必要です")
        body = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if body is not None:
            left, top, width, height = body.left, body.top, body.width, body.height
        else:
            left, top, width, height = Inches(0.5), Inches(1.8), Inches(9.0), Inches(4.5)
        add_table(slide, table_def, left, top, width, height)

    elif layout_key == "picture":
        image_path = slide_def.get("image_path")
        if not image_path:
            raise ValueError("layout が picture のスライドには image_path キーが必要です")
        image_file = Path(image_path)
        if not image_file.is_file():
            raise ValueError(f"image_path が見つかりません: {image_path}")
        body = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if body is not None:
            left, top, width = body.left, body.top, body.width
        else:
            left, top, width = Inches(1.0), Inches(1.8), Inches(8.0)
        slide.shapes.add_picture(str(image_file), left, top, width=width)
        caption = slide_def.get("caption")
        if caption:
            box = slide.shapes.add_textbox(left, prs.slide_height - Inches(0.8), width, Inches(0.5))
            box.text_frame.text = str(caption)

    notes = slide_def.get("notes")
    if notes:
        slide.notes_slide.notes_text_frame.text = str(notes)


def main() -> int:
    setup_utf8_stdio()
    parser = argparse.ArgumentParser()
    parser.add_argument("output_path")
    group = parser.add_mutually_exclusive_group(required=True)
    group.add_argument("--data")
    group.add_argument("--data-file")
    args = parser.parse_args()

    if args.data_file:
        data_file = Path(args.data_file)
        if not data_file.exists():
            print(f"データファイルが見つかりません: {args.data_file}", file=sys.stderr)
            return 1
        try:
            raw = data_file.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            raw = data_file.read_text(encoding="cp932")
    else:
        raw = args.data

    try:
        data = json.loads(raw)
    except json.JSONDecodeError as e:
        print(f"JSONとして解析できませんでした: {e}", file=sys.stderr)
        return 1

    slide_defs = data.get("slides") if isinstance(data, dict) else None
    if not slide_defs:
        print("data の 'slides' キーに1件以上のスライド定義が必要です。", file=sys.stderr)
        return 1

    if not TEMPLATE_PATH.is_file():
        print(f"テンプレートが見つかりません: {TEMPLATE_PATH}", file=sys.stderr)
        return 1
    prs = Presentation(str(TEMPLATE_PATH))
    try:
        for slide_def in slide_defs:
            apply_slide(prs, slide_def)
    except ValueError as e:
        print(f"スライド生成に失敗しました: {e}", file=sys.stderr)
        return 1

    output_path = Path(args.output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(str(output_path))
    except OSError as e:
        print(f"ファイルの保存に失敗しました: {e}", file=sys.stderr)
        return 1

    result = {
        "output_path": str(output_path),
        "total_slides": len(prs.slides),
        "size_bytes": output_path.stat().st_size,
    }
    print(json.dumps(result, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    sys.exit(main())
