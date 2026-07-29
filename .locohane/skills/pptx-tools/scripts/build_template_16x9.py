"""16:9テンプレート(assets/template_16x9.pptx)を再生成するビルドスクリプト。

python-pptxの既定テンプレートは4:3(9144000x6858000 EMU)のため、
高さを据え置いたまま幅を16:9(12192000 EMU)に拡張し、
全レイアウト・スライドマスターのプレースホルダをX方向にスケーリングする。
テンプレート仕様を見直す場合のみ実行し、生成物をassets/にコミットする。
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

OLD_WIDTH = Emu(9144000)
NEW_WIDTH = Emu(12192000)
HEIGHT = Emu(6858000)
SCALE_X = NEW_WIDTH / OLD_WIDTH

OUTPUT_PATH = Path(__file__).parent.parent / "assets" / "template_16x9.pptx"


def collect_geometry(shapes, out: list) -> None:
    # masterをスケールしてからlayoutを読むと、xfrm未定義のlayout側placeholderが
    # 継承経由でスケール済みの値を読んでしまい二重スケーリングになる。
    # そのため、全shapeの(スケール前)ジオメトリを先に読み切ってから後で一括反映する。
    for shape in shapes:
        if shape.left is not None and shape.top is not None:
            out.append((shape, shape.left, shape.top, shape.width, shape.height))


def main() -> None:
    prs = Presentation()
    prs.slide_width = NEW_WIDTH
    prs.slide_height = HEIGHT

    geometry: list = []
    for master in prs.slide_masters:
        collect_geometry(master.shapes, geometry)
        for layout in master.slide_layouts:
            collect_geometry(layout.shapes, geometry)

    for shape, left, top, width, height in geometry:
        shape.left = Emu(round(left * SCALE_X))
        shape.top = Emu(top)
        if width is not None and height is not None:
            shape.width = Emu(round(width * SCALE_X))
            shape.height = Emu(height)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    print(f"saved: {OUTPUT_PATH}")


if __name__ == "__main__":
    sys.exit(main())
